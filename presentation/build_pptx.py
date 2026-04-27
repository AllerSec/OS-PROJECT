"""
Generate the short PowerPoint deck from the captured screenshots.

Output: presentation/ObsidianLab.pptx (10 slides, 16:9, dark theme).
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Inches, Pt

HERE = Path(__file__).parent
SHOTS = HERE / "screenshots"
OUT = HERE / "ObsidianLab.pptx"

# 16:9 — 13.333" x 7.5"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Brand
BG = RGBColor(0x05, 0x06, 0x0A)
PANEL = RGBColor(0x0E, 0x12, 0x1B)
TEXT_1 = RGBColor(0xF4, 0xF6, 0xFA)
TEXT_2 = RGBColor(0xA8, 0xB1, 0xC4)
TEXT_3 = RGBColor(0x66, 0x70, 0x80)
MINT = RGBColor(0x39, 0xFF, 0xB0)
CYAN = RGBColor(0x5B, 0xE0, 0xFF)
AMBER = RGBColor(0xFF, 0xB5, 0x47)
CORAL = RGBColor(0xFF, 0x5C, 0x7A)
VIOLET = RGBColor(0x8B, 0x7C, 0xFF)


def fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def stroke(shape, color, width_pt=0.75):
    shape.line.color.rgb = color
    shape.line.width = Pt(width_pt)


def add_bg(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    fill(bg, BG)
    bg.line.fill.background()
    return bg


def add_text(slide, text, left, top, width, height, *,
             size=18, color=TEXT_1, bold=False, font="Calibri",
             align="left"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(0); tf.margin_right = Pt(0)
    tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font
    if align == "center":
        from pptx.enum.text import PP_ALIGN
        p.alignment = PP_ALIGN.CENTER
    return box


def add_glass_panel(slide, left, top, width, height):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    rect.adjustments[0] = 0.04
    rect.fill.solid()
    rect.fill.fore_color.rgb = PANEL
    rect.line.color.rgb = RGBColor(0x22, 0x28, 0x36)
    rect.line.width = Pt(0.75)
    return rect


def add_accent_line(slide, top, color=MINT):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.6), top,
                                   Inches(0.3), Inches(0.04))
    fill(line, color)


def add_pill(slide, text, left, top, color):
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, top, Inches(1.6), Inches(0.36))
    pill.adjustments[0] = 0.5
    pill.fill.solid()
    pill.fill.fore_color.rgb = PANEL
    pill.line.color.rgb = color
    pill.line.width = Pt(1)
    tf = pill.text_frame
    tf.margin_left = Pt(8); tf.margin_right = Pt(8)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(11)
    p.font.color.rgb = color
    p.font.bold = True
    p.font.name = "Consolas"
    from pptx.enum.text import PP_ALIGN
    p.alignment = PP_ALIGN.CENTER


def insert_image(slide, path, left, top, width=None, height=None):
    if not Path(path).exists():
        return None
    if width and not height:
        return slide.shapes.add_picture(str(path), left, top, width=width)
    if height and not width:
        return slide.shapes.add_picture(str(path), left, top, height=height)
    if width and height:
        return slide.shapes.add_picture(str(path), left, top, width=width, height=height)
    return slide.shapes.add_picture(str(path), left, top)


def slide_title(prs, title, subtitle, kicker="OBSIDIAN LAB"):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)

    # Subtle accent bar at top
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.06))
    fill(bar, MINT)

    add_text(s, kicker, Inches(0.6), Inches(0.45),
             Inches(8), Inches(0.4), size=12, color=CYAN, bold=True, font="Consolas")

    add_text(s, title, Inches(0.6), Inches(2.1),
             Inches(12), Inches(2.0), size=54, color=TEXT_1, bold=True)

    add_text(s, subtitle, Inches(0.6), Inches(4.4),
             Inches(11), Inches(1.8), size=20, color=TEXT_2)

    add_text(s, "Real-Time Operating Systems · 2025–2026", Inches(0.6),
             Inches(6.7), Inches(8), Inches(0.4), size=12, color=TEXT_3,
             font="Consolas")
    return s


def slide_with_image(prs, kicker, title, subtitle, img_path, *,
                     img_w_in=8.5, img_left_in=4.4, img_top_in=1.05,
                     bullets=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)

    add_accent_line(s, Inches(0.55), color=MINT)
    add_text(s, kicker, Inches(0.6), Inches(0.7),
             Inches(8), Inches(0.4), size=12, color=CYAN, bold=True, font="Consolas")
    add_text(s, title, Inches(0.6), Inches(1.05),
             Inches(4.0), Inches(1.5), size=30, color=TEXT_1, bold=True)
    add_text(s, subtitle, Inches(0.6), Inches(2.45),
             Inches(3.6), Inches(2.5), size=14, color=TEXT_2)

    if bullets:
        y = 3.85
        for b in bullets:
            add_text(s, "•  " + b, Inches(0.6), Inches(y),
                     Inches(3.7), Inches(0.5), size=12, color=TEXT_2)
            y += 0.45

    if img_path and Path(img_path).exists():
        insert_image(s, img_path, Inches(img_left_in), Inches(img_top_in),
                     width=Inches(img_w_in))

    return s


def slide_section(prs, kicker, title, body):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_accent_line(s, Inches(0.55), color=CYAN)
    add_text(s, kicker, Inches(0.6), Inches(0.7),
             Inches(8), Inches(0.4), size=12, color=CYAN, bold=True, font="Consolas")
    add_text(s, title, Inches(0.6), Inches(1.05),
             Inches(12), Inches(1.6), size=38, color=TEXT_1, bold=True)
    add_text(s, body, Inches(0.6), Inches(2.7),
             Inches(11.5), Inches(4.2), size=18, color=TEXT_2)
    return s


def slide_six_demos(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_accent_line(s, Inches(0.55), color=AMBER)
    add_text(s, "OVERVIEW", Inches(0.6), Inches(0.7),
             Inches(8), Inches(0.4), size=12, color=CYAN, bold=True, font="Consolas")
    add_text(s, "Six interactive demos, one page", Inches(0.6), Inches(1.05),
             Inches(12), Inches(1.0), size=30, color=TEXT_1, bold=True)

    demos = [
        ("01", "Race condition", "Two threads, one counter — with and without a mutex.", MINT),
        ("02", "Mutex vs spinlock", "Same workload, two locks: sleep vs busy-wait.", CYAN),
        ("03", "Producer / consumer", "Bounded buffer with two semaphores + mutex.", AMBER),
        ("04", "Readers / writers", "Many readers may share access; writers go alone.", VIOLET),
        ("05", "Deadlock", "Circular wait — and the global lock-order fix.", CORAL),
        ("06", "Priority inversion", "The bug + the priority-inheritance solution.", MINT),
    ]

    grid_left = Inches(0.6)
    grid_top = Inches(2.2)
    cw = Inches(4.05)
    ch = Inches(1.55)
    gx = Inches(0.18)
    gy = Inches(0.18)

    for i, (num, title, sub, color) in enumerate(demos):
        col = i % 3
        row = i // 3
        x = grid_left + (cw + gx) * col
        y = grid_top + (ch + gy) * row
        panel = add_glass_panel(s, x, y, cw, ch)
        # number badge
        add_text(s, num, x + Inches(0.2), y + Inches(0.18),
                 Inches(0.6), Inches(0.3), size=11, color=color, bold=True, font="Consolas")
        # title
        add_text(s, title, x + Inches(0.2), y + Inches(0.45),
                 Inches(3.7), Inches(0.45), size=18, color=TEXT_1, bold=True)
        # subtitle
        add_text(s, sub, x + Inches(0.2), y + Inches(0.92),
                 Inches(3.7), Inches(0.65), size=11, color=TEXT_2)

    return s


def slide_architecture(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_accent_line(s, Inches(0.55), color=VIOLET)
    add_text(s, "ARCHITECTURE", Inches(0.6), Inches(0.7),
             Inches(8), Inches(0.4), size=12, color=CYAN, bold=True, font="Consolas")
    add_text(s, "How it fits together", Inches(0.6), Inches(1.05),
             Inches(12), Inches(1.0), size=30, color=TEXT_1, bold=True)

    # Browser box
    add_glass_panel(s, Inches(0.8), Inches(2.5), Inches(4.2), Inches(2.6))
    add_text(s, "BROWSER", Inches(1.0), Inches(2.65),
             Inches(3.5), Inches(0.4), size=11, color=CYAN, bold=True, font="Consolas")
    add_text(s, "Vanilla JS\nChart.js\nLiquid-glass CSS",
             Inches(1.0), Inches(3.05), Inches(4.0), Inches(2.0),
             size=15, color=TEXT_1)

    # Server box
    add_glass_panel(s, Inches(8.2), Inches(2.5), Inches(4.2), Inches(2.6))
    add_text(s, "SERVER", Inches(8.4), Inches(2.65),
             Inches(3.5), Inches(0.4), size=11, color=MINT, bold=True, font="Consolas")
    add_text(s, "FastAPI + uvicorn\nthreading + sync primitives\n6 demo modules",
             Inches(8.4), Inches(3.05), Inches(4.0), Inches(2.0),
             size=15, color=TEXT_1)

    # Arrow
    arrow = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.1), Inches(3.65),
                                 Inches(3.0), Inches(0.06))
    fill(arrow, CYAN)
    add_text(s, "WebSocket  /ws/{demo}", Inches(5.1), Inches(3.05),
             Inches(3.0), Inches(0.45), size=12, color=CYAN, bold=True,
             font="Consolas", align="center")

    add_text(s,
             "One WebSocket route per demo. Each demo spawns real OS threads, "
             "uses real Linux-style sync primitives (mutex, semaphore, RW-lock), "
             "and pushes events live to the browser as JSON.",
             Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.6),
             size=14, color=TEXT_2)
    return s


def slide_run(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    add_accent_line(s, Inches(0.55), color=MINT)
    add_text(s, "RUN IT", Inches(0.6), Inches(0.7),
             Inches(8), Inches(0.4), size=12, color=CYAN, bold=True, font="Consolas")
    add_text(s, "Three commands", Inches(0.6), Inches(1.05),
             Inches(12), Inches(1.0), size=30, color=TEXT_1, bold=True)

    code = (
        "$ pip install -r app/requirements.txt\n\n"
        "$ uvicorn app.main:app --reload\n\n"
        "$ open http://localhost:8000"
    )
    panel = add_glass_panel(s, Inches(0.8), Inches(2.4), Inches(11.7), Inches(3.4))
    add_text(s, code, Inches(1.1), Inches(2.7), Inches(11.0), Inches(3.0),
             size=22, color=MINT, font="Consolas", bold=True)

    add_text(s, "No Node, no build step, no database. Works on Windows, macOS, and Linux.",
             Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.6),
             size=14, color=TEXT_2)
    return s


def slide_closing(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)

    add_text(s, "Thank you.", Inches(0.6), Inches(2.7),
             Inches(12), Inches(2), size=72, color=TEXT_1, bold=True)
    add_text(s, "Questions?", Inches(0.6), Inches(4.2),
             Inches(12), Inches(1), size=24, color=CYAN, font="Consolas")

    # Repo + live demo links
    add_text(s, "Repository", Inches(0.6), Inches(5.8),
             Inches(4), Inches(0.3), size=11, color=TEXT_3, bold=True, font="Consolas")
    add_text(s, "github.com/AllerSec/OS-PROJECT", Inches(0.6), Inches(6.1),
             Inches(8), Inches(0.5), size=18, color=MINT, font="Consolas", bold=True)

    add_text(s, "Live demo", Inches(0.6), Inches(6.7),
             Inches(4), Inches(0.3), size=11, color=TEXT_3, bold=True, font="Consolas")
    add_text(s, "allersec.github.io/OS-PROJECT", Inches(0.6), Inches(7.0),
             Inches(8), Inches(0.5), size=18, color=CYAN, font="Consolas", bold=True)
    return s


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 1 — Title
    slide_title(prs,
                "Linux Synchronization Tools for RTOS",
                "An interactive web app demonstrating mutex, spinlock, semaphores,\n"
                "RW-lock, deadlock and priority inversion — live in the browser.")

    # 2 — Why this matters
    slide_section(prs, "MOTIVATION",
                  "Why synchronization in RTOS?",
                  "In real-time systems, threads share buffers, devices and counters. "
                  "When two of them touch the same data without coordination, the result "
                  "is silently wrong, the system deadlocks, or a high-priority task misses "
                  "its deadline.\n\nLinux provides the tools to prevent this — but each "
                  "comes with trade-offs. This project demonstrates them, side by side, "
                  "with code you can run.")

    # 3 — The six demos at a glance
    slide_six_demos(prs)

    # 4 — Architecture
    slide_architecture(prs)

    # 5 — UI hero
    slide_with_image(prs,
                     "DESIGN",
                     "Liquid-glass dark UI",
                     "Single-page app, fully responsive\nfor desktop and mobile.\n\n"
                     "Tailwind-inspired tokens, custom CSS\n"
                     "(no framework), Geist + JetBrains Mono.",
                     SHOTS / "01_hero.png",
                     img_w_in=8.4, img_left_in=4.5, img_top_in=2.4)

    # 6 — Race condition (with results)
    slide_with_image(prs,
                     "DEMO 01",
                     "Race condition",
                     "Two threads each increment\na shared counter.\n\n"
                     "Without a mutex, updates are lost.\n"
                     "With a mutex, the count is exact.",
                     SHOTS / "08_race_condition_results.png",
                     img_w_in=8.0, img_left_in=4.7, img_top_in=1.05,
                     bullets=["Real reproduction (not simulated)",
                              "Side-by-side: unsafe vs safe",
                              "Lost updates counted live"])

    # 7 — Mutex vs spinlock
    slide_with_image(prs,
                     "DEMO 02",
                     "Mutex vs spinlock",
                     "Same workload run with both.\n\n"
                     "Mutex sleeps when blocked.\n"
                     "Spinlock burns CPU until free.\n\n"
                     "Linux kernel uses both, depending\non the context.",
                     SHOTS / "09_mvs_results.png",
                     img_w_in=8.0, img_left_in=4.7, img_top_in=1.05)

    # 8 — Producer / consumer running
    slide_with_image(prs,
                     "DEMO 03",
                     "Producer / consumer",
                     "Bounded buffer of 8 slots.\n\n"
                     "Two semaphores (empty + full)\nplus a mutex protect the buffer.\n\n"
                     "The animation tracks the real\nbuffer state, slot by slot.",
                     SHOTS / "11_producer_consumer_running.png",
                     img_w_in=8.4, img_left_in=4.5, img_top_in=1.05)

    # 9 — Deadlock
    slide_with_image(prs,
                     "DEMO 05",
                     "Deadlock",
                     "Two threads grab two locks in\nopposite order — they wait forever.\n\n"
                     "A 2.5-second watchdog gives up\nand reports DEADLOCKED.\n\n"
                     'Click "Fix" to apply a global\nlock order. The cycle disappears.',
                     SHOTS / "10_deadlock_broken.png",
                     img_w_in=8.0, img_left_in=4.7, img_top_in=1.05)

    # 10 — Mobile
    slide_with_image(prs,
                     "RESPONSIVE",
                     "Designed for mobile too",
                     "All six demos stack into a single\ncolumn on small screens.\n\n"
                     "Reduced blur for performance.\n"
                     "Touch targets ≥ 44 px.\n"
                     "Respects prefers-reduced-motion.",
                     SHOTS / "12_mobile_view.png",
                     img_w_in=2.5, img_left_in=10.3, img_top_in=1.05)

    # 11 — How to run
    slide_run(prs)

    # 12 — Closing
    slide_closing(prs)

    OUT.parent.mkdir(exist_ok=True)
    prs.save(str(OUT))
    print(f"wrote {OUT}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
